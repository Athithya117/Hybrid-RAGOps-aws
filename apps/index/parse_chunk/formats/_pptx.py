#!/usr/bin/env python3
from __future__ import annotations
import os
import sys
import json
import time
import hashlib
import tempfile
import re
import unicodedata
import threading
import io
from io import BytesIO
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple

class LoggerShim:
    def __init__(self, name: str): self.name = name
    def _emit(self, level: str, event: str, msg: str = "", **extra):
        o = {"ts": datetime.utcnow().isoformat() + "Z", "level": level, "event": event, "msg": msg}
        if extra: o.update(extra)
        print(json.dumps(o, ensure_ascii=False), flush=True)
    def _unpack(self, a, b, fmt_args, kwargs, default_event):
        if b is None:
            event = kwargs.pop("event", default_event); msg = a
        else:
            event = a; msg = b
        if fmt_args:
            try: msg = msg % fmt_args
            except Exception:
                try: msg = msg.format(*fmt_args)
                except Exception: pass
        return event, msg, kwargs
    def info(self, a, b=None, *fmt_args, **kwargs): e,m,k=self._unpack(a,b,fmt_args,kwargs,"info"); self._emit("info",e,m,**k)
    def warning(self, a, b=None, *fmt_args, **kwargs): e,m,k=self._unpack(a,b,fmt_args,kwargs,"warn"); self._emit("warn",e,m,**k)
    def warn(self, a, b=None, *fmt_args, **kwargs): self.warning(a,b,*fmt_args,**kwargs)
    def error(self, a, b=None, *fmt_args, **kwargs): e,m,k=self._unpack(a,b,fmt_args,kwargs,"error"); self._emit("error",e,m,**k)
    def exception(self, a, b=None, *fmt_args, **kwargs):
        import traceback
        tb = traceback.format_exc()
        e,m,k=self._unpack(a,b,fmt_args,kwargs,"exception"); k.update({"tb": tb}); self._emit("error",e,m,**k)

log = LoggerShim("pptx_parser")

USE_MANAGED_IDENTITY = os.getenv("AZURE_USE_MANAGED_IDENTITY", "").strip().lower() in ("1", "true", "yes")
ENV = os.getenv("ENV", "STAGING").upper()
STORAGE_PROTOCOL = "az"
AZURE_CONTAINER: Optional[str] = os.getenv("AZURE_CONTAINER") or os.getenv("STORAGE_CONTAINER") or os.getenv("AZ_CONTAINER")
AZURE_RAW_PREFIX: str = os.getenv("AZURE_RAW_PREFIX", os.getenv("STORAGE_RAW_PREFIX", "data/raw/")).rstrip("/") + "/"
AZURE_CHUNKED_PREFIX: str = os.getenv("AZURE_CHUNKED_PREFIX", os.getenv("STORAGE_CHUNKED_PREFIX", "data/chunked/")).rstrip("/") + "/"
SLIDES_PER_CHUNK: int = int(os.getenv("PPTX_SLIDES_PER_CHUNK", "3"))
DISABLE_OCR: bool = os.getenv("PPTX_DISABLE_OCR", "false").lower() == "true"
FORCE_OCR: bool = os.getenv("PPTX_FORCE_OCR", "false").lower() == "true"
OCR_BACKEND: str = os.getenv("PPTX_OCR_ENGINE", "tesseract").lower()
PPTX_OCR_STRICT: bool = os.getenv("PPTX_OCR_STRICT", "false").lower() == "true"
MIN_IMG_BYTES: int = int(os.getenv("PPTX_MIN_IMG_SIZE_BYTES", "3072"))
PARSER_VERSION_PPTX: str = os.getenv("PARSER_VERSION_PPTX", "pptx-parser-v1")
TOKEN_ENCODER: str = os.getenv("TOKEN_ENCODER", "cl100k_base")
FORCE_OVERWRITE: bool = os.getenv("FORCE_OVERWRITE", "false").lower() == "true"
CHUNKED_SCHEMA_VERSION: str = os.getenv("CHUNKED_SCHEMA_VERSION", "chunked_v1")
PUT_RETRIES = int(os.getenv("PUT_RETRIES", "3"))
PUT_BACKOFF = float(os.getenv("PUT_BACKOFF", "0.3"))

_fs = None
_pa = None
_pq = None
_np = None
_storage_root = None
_storage_lock = threading.Lock()
_storage_client = None
BLOB_CLIENT = None
AZURE_SDK_LOADED = False

try:
    import logging as _lg
    _root = _lg.getLogger()
    _root.setLevel(_lg.WARNING)
    for _n in ("adlfs", "azure", "azure.storage", "azure.core", "azure.identity", "urllib3", "botocore", "requests", "httpx"):
        lg = _lg.getLogger(_n)
        lg.setLevel(_lg.WARNING)
        lg.propagate = False
except Exception:
    pass

def _init_env():
    global _fs, _pa, _pq, _np, _storage_root, BLOB_CLIENT, AZURE_SDK_LOADED
    if getattr(_init_env, "done", False): return
    if not AZURE_CONTAINER:
        log.error("startup_missing_container", "AZURE_CONTAINER (or STORAGE_CONTAINER/AZ_CONTAINER) must be set")
        sys.exit(1)
    if USE_MANAGED_IDENTITY:
        try:
            from azure.identity import DefaultAzureCredential
            from azure.storage.blob import BlobServiceClient
            AZURE_SDK_LOADED = True
        except Exception as e:
            log.error("azure_sdk_missing", "AZURE_USE_MANAGED_IDENTITY requested but azure-identity/azure-storage-blob not installed", error=str(e))
            sys.exit(2)
        account_name = os.environ.get("AZURE_STORAGE_ACCOUNT_NAME") or os.environ.get("AZURE_ACCOUNT_NAME")
        if not account_name:
            log.error("missing_account_name", "AZURE_STORAGE_ACCOUNT_NAME required for managed identity mode")
            sys.exit(2)
        endpoint_suffix = os.environ.get("AZURE_ENDPOINT_SUFFIX", "core.windows.net")
        account_url = f"https://{account_name}.blob.{endpoint_suffix}"
        try:
            uai_client_id = os.environ.get("UAI_RAG_RW_CLIENT_ID") or os.environ.get("AZURE_CLIENT_ID") or None
            if uai_client_id:
                cred = DefaultAzureCredential(managed_identity_client_id=uai_client_id)
            else:
                cred = DefaultAzureCredential()
            BLOB_CLIENT = BlobServiceClient(account_url=account_url, credential=cred, connection_timeout=60)
            _fs = None
            _storage_root = f"az://{AZURE_CONTAINER.rstrip('/')}/"
            try:
                container_client = BLOB_CLIENT.get_container_client(AZURE_CONTAINER)
                try:
                    container_client.get_container_properties()
                    log.info("azure_mi_ready", "Managed identity BlobServiceClient initialized and container validated", account=account_name)
                except Exception as e:
                    log.warning("mi_smoke", "managed identity client created, but smoke-check failed (may be normal in restricted env)", error=str(e))
            except Exception:
                log.info("azure_mi_ready", "Managed identity BlobServiceClient initialized", account=account_name)
        except Exception as e:
            log.error("blob_client_init_failed", "Failed to create BlobServiceClient with Managed Identity", error=str(e))
            raise
    else:
        conn = os.environ.get("AZURE_STORAGE_CONNECTION_STRING")
        opts = {}
        if conn:
            opts["connection_string"] = conn
        else:
            acct = os.environ.get("AZURE_STORAGE_ACCOUNT_NAME") or os.environ.get("AZURE_ACCOUNT_NAME")
            key = os.environ.get("AZURE_STORAGE_ACCOUNT_KEY") or os.environ.get("AZURE_ACCOUNT_KEY")
            sas = os.environ.get("AZURE_SAS_TOKEN")
            eps = os.environ.get("AZURE_ENDPOINT_SUFFIX") or "core.windows.net"
            if acct and key:
                opts["account_name"] = acct; opts["account_key"] = key; opts["endpoint_suffix"] = eps
            elif acct and sas:
                opts["account_name"] = acct; opts["sas_token"] = sas; opts["endpoint_suffix"] = eps
            elif os.environ.get("AZURE_ANON"):
                if acct: opts["account_name"] = acct
                opts["anon"] = True
        if not opts:
            log.error("non_mi_no_creds", "Non-managed identity mode selected but no connection_string/account_key/SAS/anon present")
            sys.exit(2)
        try:
            import fsspec as _fsspec
            _fs = _fsspec.filesystem("az", **opts)
            _storage_root = f"az://{AZURE_CONTAINER.rstrip('/')}/"
            log.info("adlfs_ready", "adlfs filesystem initialized for Azure storage", opts_keys=list(opts.keys()))
        except Exception as e:
            log.error("adlfs_init_failed", "Cannot initialize adlfs/fsspec filesystem", error=str(e))
            raise
    try:
        import numpy as _n; _np = _n
    except Exception:
        _np = None
    try:
        import pyarrow as pa; import pyarrow.parquet as pq; _pa = pa; _pq = pq
    except Exception:
        _pa = None; _pq = None
    _init_env.done = True

def full_path_from_key(key: str) -> str:
    return (_storage_root + key.lstrip("/")) if _storage_root else key

def strip_root_from_path(full: str) -> str:
    if _storage_root and full.startswith(_storage_root): return full[len(_storage_root):]
    proto_prefix = "az://"
    if full.startswith(proto_prefix):
        rest = full[len(proto_prefix):]
        if rest.startswith(AZURE_CONTAINER + "/"): return rest[len(AZURE_CONTAINER) + 1:]
        if rest == AZURE_CONTAINER: return ""
    if full.startswith(AZURE_CONTAINER + "/"): return full[len(AZURE_CONTAINER) + 1:]
    return full

def retry(func, retries: int = 3, delay: float = 1.0, backoff: float = 2.0):
    for attempt in range(retries):
        try: return func()
        except Exception as e:
            if attempt == retries - 1: raise
            log.warning("retry_attempt", "attempt=%d error=%s", attempt + 1, str(e)); time.sleep(delay); delay *= backoff

class AzureStorageClient:
    def __init__(self, fs_obj=None, root=None, container=None, blob_client=None):
        self.fs = fs_obj
        self.root = root
        self.container = container
        self.blob_client = blob_client
    def _container_client(self):
        if self.blob_client is None:
            raise RuntimeError("blob_client not initialized for managed-identity mode")
        return self.blob_client.get_container_client(self.container)
    def head_object(self, Bucket, Key):
        if self.fs is not None:
            full = full_path_from_key(Key)
            info = self.fs.info(full)
            out = {"ContentLength": int(info.get("size", 0))}
            etag = info.get("etag") or info.get("ETag") or info.get("eTag") or ""
            out["ETag"] = etag
            out["LastModified"] = info.get("Last-Modified") or info.get("last_modified") or info.get("LastModified") or ""
            out["Metadata"] = info.get("metadata") or info.get("meta") or {}
            return out
        else:
            container_client = self._container_client()
            blob_client = container_client.get_blob_client(Key)
            props = blob_client.get_blob_properties()
            out = {
                "ContentLength": int(getattr(props, "size", 0) or 0),
                "ETag": getattr(props, "etag", "") or "",
                "LastModified": getattr(props, "last_modified", "") or "",
                "Metadata": getattr(props, "metadata", {}) or {},
            }
            return out
    def get_object(self, Bucket, Key):
        if self.fs is not None:
            full = full_path_from_key(Key)
            with self.fs.open(full, "rb") as f:
                data = f.read()
            return {"Body": BytesIO(data)}
        else:
            container_client = self._container_client()
            blob_client = container_client.get_blob_client(Key)
            stream = blob_client.download_blob()
            data = stream.readall()
            return {"Body": BytesIO(data)}
    def put_object(self, Bucket, Key, Body, ContentType=None):
        if self.fs is not None:
            full = full_path_from_key(Key)
            if isinstance(Body, (bytes, bytearray)):
                b = Body
            elif isinstance(Body, str):
                b = Body.encode("utf-8")
            elif hasattr(Body, "read"):
                b = Body.read()
                if isinstance(b, str):
                    b = b.encode("utf-8")
            else:
                try: b = bytes(Body)
                except Exception: b = str(Body).encode("utf-8")
            with self.fs.open(full, "wb") as f: f.write(b)
            return {"ResponseMetadata": {"HTTPStatusCode": 200}}
        else:
            container_client = self._container_client()
            blob_client = container_client.get_blob_client(Key)
            if isinstance(Body, (bytes, bytearray)):
                data = Body
            elif isinstance(Body, str):
                data = Body.encode("utf-8")
            elif hasattr(Body, "read"):
                data = Body.read()
                if isinstance(data, str):
                    data = data.encode("utf-8")
            else:
                data = str(Body).encode("utf-8")
            blob_client.upload_blob(data, overwrite=True)
            return {"ResponseMetadata": {"HTTPStatusCode": 200}}
    def upload_file(self, LocalFile, Bucket, Key, ExtraArgs=None):
        if self.fs is not None:
            full = full_path_from_key(Key)
            if hasattr(self.fs, "put"):
                self.fs.put(LocalFile, full)
            else:
                with open(LocalFile, "rb") as lf: d = lf.read()
                with self.fs.open(full, "wb") as f: f.write(d)
            return
        else:
            container_client = self._container_client()
            blob_client = container_client.get_blob_client(Key)
            with open(LocalFile, "rb") as lf:
                blob_client.upload_blob(lf, overwrite=True)
    def copy_object(self, CopySource, Bucket, Key):
        src = CopySource.get("Key")
        if self.fs is not None:
            full_src = full_path_from_key(src); full_dst = full_path_from_key(Key)
            with self.fs.open(full_src, "rb") as rf: data = rf.read()
            with self.fs.open(full_dst, "wb") as wf: wf.write(data)
            return
        else:
            src_blob_client = self._container_client().get_blob_client(src)
            dst_blob_client = self._container_client().get_blob_client(Key)
            src_url = src_blob_client.url
            dst_blob_client.start_copy_from_url(src_url)
    def delete_object(self, Bucket, Key):
        if self.fs is not None:
            full = full_path_from_key(Key)
            try: self.fs.rm(full)
            except Exception:
                try: self.fs.delete(full)
                except Exception: pass
            return
        else:
            blob_client = self._container_client().get_blob_client(Key)
            try: blob_client.delete_blob()
            except Exception: pass
    def exists(self, full_key: str) -> bool:
        if self.fs is not None:
            try:
                return self.fs.exists(full_path_from_key(full_key))
            except Exception:
                return False
        else:
            container_client = self._container_client()
            blobs = list(container_client.list_blobs(name_starts_with=full_key))
            return len(blobs) > 0
    def get_paginator(self, name):
        if self.fs is not None:
            class P:
                def __init__(self, fs, root): self.fs = fs; self.root = root
                def paginate(self, Bucket, Prefix, PaginationConfig=None):
                    base = (Prefix.rstrip("/")) + "/"; root_path = self.root + base
                    try:
                        if hasattr(self.fs, "find"): found = self.fs.find(root_path)
                        else: found = self.fs.glob(root_path + "**", recursive=True)
                    except Exception: found = []
                    page = {"Contents": []}
                    for f in found:
                        try: info = self.fs.info(f)
                        except Exception: continue
                        if info.get("type") == "directory": continue
                        rel = strip_root_from_path(f); page["Contents"].append({"Key": rel})
                        if len(page["Contents"]) >= 1000: yield page; page = {"Contents": []}
                    if page["Contents"]: yield page
            return P(self.fs, self.root)
        else:
            class Pblob:
                def __init__(self, container_client): self.container_client = container_client
                def paginate(self, Bucket, Prefix, PaginationConfig=None):
                    blobs = self.container_client.list_blobs(name_starts_with=Prefix)
                    page = {"Contents": []}
                    for b in blobs:
                        page["Contents"].append({"Key": b.name})
                        if len(page["Contents"]) >= 1000:
                            yield page
                            page = {"Contents": []}
                    if page["Contents"]: yield page
            return Pblob(self._container_client())

def get_storage_client_singleton():
    global _storage_client
    if _storage_client is None:
        with _storage_lock:
            if _storage_client is None:
                _init_env()
                if USE_MANAGED_IDENTITY:
                    _storage_client = AzureStorageClient(fs_obj=None, root=_storage_root, container=AZURE_CONTAINER, blob_client=BLOB_CLIENT)
                else:
                    _storage_client = AzureStorageClient(fs_obj=_fs, root=_storage_root, container=AZURE_CONTAINER, blob_client=None)
    return _storage_client

def sha256_hex_bytes(b: bytes) -> str: return hashlib.sha256(b).hexdigest()
def sha256_hex_str(s: str) -> str: return hashlib.sha256((s or "").encode("utf-8")).hexdigest()
def canonicalize_text(s: str) -> str:
    if not isinstance(s, str): s = str(s or "")
    s = unicodedata.normalize("NFKC", s); s = s.replace("\r\n", "\n").replace("\r", "\n")
    lines = [re.sub(r"[ \t]+$", "", ln) for ln in s.split("\n")]
    return "\n".join(lines).strip()

def _load_tiktoken_encoder(name: str):
    try:
        import tiktoken
        enc = None
        try: enc = tiktoken.encoding_for_model(name)
        except Exception:
            try: enc = tiktoken.get_encoding(name)
            except Exception: enc = None
        return enc
    except Exception: return None

def _count_tokens(text: str) -> int:
    if not text: return 0
    enc = _load_tiktoken_encoder(TOKEN_ENCODER)
    if enc:
        try: return len(enc.encode(text))
        except Exception: pass
    return len(text.split())

def is_ocr_line_valid(text: str, min_ratio: float = 0.6) -> bool:
    t = (text or "").strip()
    if len(t) < 5: return False
    alnum = sum(c.isalnum() for c in t)
    try: return (alnum / len(t)) >= min_ratio
    except Exception: return False

def dedupe_lines(lines: list) -> list:
    seen, out = set(), []
    for l in lines:
        key = (l or "").strip().lower()
        if key and key not in seen: seen.add(key); out.append(l)
    return out

def storage_upload_file_atomic(local_path: str, key: str, content_type: str = "application/octet-stream"):
    _init_env()
    client = get_storage_client_singleton()
    full = full_path_from_key(key)
    tmp = f"{full}.tmp.{os.getpid()}.{int(time.time())}"
    for attempt in range(1, PUT_RETRIES + 1):
        try:
            if client.fs is not None and hasattr(client.fs, "put"):
                client.fs.put(local_path, tmp)
            elif client.fs is not None:
                with open(local_path, "rb") as lf: d = lf.read()
                with client.fs.open(tmp, "wb") as f: f.write(d)
            else:
                client.upload_file(local_path, AZURE_CONTAINER, key)
                return
            if client.fs is not None and hasattr(client.fs, "mv"):
                client.fs.mv(tmp, full)
            else:
                with client.fs.open(tmp, "rb") as rf: data = rf.read()
                with client.fs.open(full, "wb") as wf: wf.write(data)
                try: client.fs.rm(tmp)
                except Exception: pass
            return
        except Exception as e:
            log.warning("upload_retry", "attempt=%d key=%s error=%s", attempt, key, str(e)); time.sleep(PUT_BACKOFF * attempt)
    raise Exception(f"atomic upload failed for {key} after {PUT_RETRIES} attempts")

class ParquetWriter:
    def __init__(self, doc_id: str): self.doc_id = doc_id; self._rows: List[Dict[str, Any]] = []
    def _normalize(self, payload: Dict[str, Any]) -> Dict[str, Any]:
        fields: Dict[str, Any] = {}
        fields["document_id"] = payload.get("document_id") or ""
        fields["file_name"] = payload.get("file_name") or ""
        fields["chunk_id"] = payload.get("chunk_id") or ""
        fields["chunk_type"] = payload.get("chunk_type") or ""
        fields["text"] = canonicalize_text(payload.get("text") or "")
        try: fields["token_count"] = int(payload.get("token_count") or 0)
        except Exception: fields["token_count"] = 0
        for k in ("figures", "tags", "layout_tags", "heading_path", "headings"):
            v = payload.get(k, None)
            try:
                if v is None: fields[k] = "[]"
                else: fields[k] = json.dumps(v, ensure_ascii=False, sort_keys=True)
            except Exception:
                fields[k] = "[]"
        fields["file_type"] = payload.get("file_type") or ""
        fields["source_url"] = payload.get("source_url") or ""
        fields["slide_start"] = None; fields["slide_end"] = None
        if payload.get("slide_range") and isinstance(payload.get("slide_range"), (list, tuple)) and len(payload.get("slide_range")) >= 2:
            try: fields["slide_start"] = int(payload["slide_range"][0]); fields["slide_end"] = int(payload["slide_range"][1])
            except Exception: fields["slide_start"] = None; fields["slide_end"] = None
        fields["timestamp"] = payload.get("timestamp") or ""
        fields["parser_version"] = payload.get("parser_version") or PARSER_VERSION_PPTX
        fields["used_ocr"] = bool(payload.get("used_ocr", False)); fields["layout"] = payload.get("layout") or ""
        return fields
    def write_payload(self, payload: Dict[str, Any]) -> int:
        self._rows.append(self._normalize(payload)); return 1
    def finalize_and_upload(self, out_basename: str) -> Tuple[int, str, str, int]:
        _init_env()
        if _pa is None or _pq is None: raise RuntimeError("pyarrow is required to write parquet output")
        if not self._rows: return 0, "", "", 0
        schema = _pa.schema([
            _pa.field("document_id", _pa.string()), _pa.field("file_name", _pa.string()), _pa.field("chunk_id", _pa.string()),
            _pa.field("chunk_type", _pa.string()), _pa.field("text", _pa.string()), _pa.field("token_count", _pa.int64()),
            _pa.field("figures", _pa.string()), _pa.field("tags", _pa.string()), _pa.field("layout_tags", _pa.string()),
            _pa.field("heading_path", _pa.string()), _pa.field("headings", _pa.string()), _pa.field("file_type", _pa.string()),
            _pa.field("source_url", _pa.string()), _pa.field("slide_start", _pa.int64()), _pa.field("slide_end", _pa.int64()),
            _pa.field("timestamp", _pa.string()), _pa.field("parser_version", _pa.string()), _pa.field("used_ocr", _pa.bool_()), _pa.field("layout", _pa.string())
        ])
        cols = {name: [] for name in [f.name for f in schema]}
        for r in self._rows:
            for name in cols: cols[name].append(r.get(name) if name in r else None)
        table = _pa.Table.from_pydict(cols, schema=schema)
        existing_md = table.schema.metadata or {}
        new_md = dict(existing_md); new_md.update({b"schema_version": CHUNKED_SCHEMA_VERSION.encode("utf-8"), b"parser_version": PARSER_VERSION_PPTX.encode("utf-8"), b"producer": b"pptx_parser", b"created_at": datetime.utcnow().isoformat().encode("utf-8")})
        table = table.replace_schema_metadata(new_md)
        tmpfile = tempfile.NamedTemporaryFile(mode="wb", delete=False, suffix=".parquet", dir="/tmp"); tmpfile.close()
        _pq.write_table(table, tmpfile.name, compression="zstd", flavor="spark")
        local_parquet_path = tmpfile.name
        with open(local_parquet_path, "rb") as fh: b = fh.read()
        sha = sha256_hex_bytes(b); size = os.path.getsize(local_parquet_path)
        parquet_key = out_basename + ".parquet"
        storage_upload_file_atomic(local_parquet_path, AZURE_CHUNKED_PREFIX + parquet_key, content_type="application/octet-stream")
        try: os.unlink(local_parquet_path)
        except Exception: pass
        return len(self._rows), AZURE_CHUNKED_PREFIX + parquet_key, sha, size

def sanitize_payload_for_raw_manifest(doc_id: str, raw_key: str, chunked_key: str, rows: int, sha: str, size: int) -> Dict[str, Any]:
    return {"raw_key": raw_key, "doc_id": doc_id, "chunked_key": chunked_key, "rows": rows, "sha256": sha, "size_bytes": size, "schema_version": CHUNKED_SCHEMA_VERSION, "parser_version": PARSER_VERSION_PPTX, "created_at": datetime.utcnow().isoformat() + "Z"}

def _extract_image_blob_from_shape(shape):
    try:
        img = getattr(shape, "image", None)
        if img and getattr(img, "blob", None): return img.blob
    except Exception: pass
    try:
        fill = getattr(shape, "fill", None)
        if fill is not None and getattr(fill, "type", None) is not None:
            pic = getattr(fill, "picture", None)
            if pic and getattr(pic, "image", None) and getattr(pic.image, "blob", None): return pic.image.blob
    except Exception: pass
    return None

def do_ocr(img: Any) -> list:
    lines = []
    try:
        if OCR_BACKEND == "tesseract":
            try:
                import cv2; from PIL import Image; import pytesseract
            except Exception as e:
                log.warning("ocr_import_failed", "Tesseract backend requested but imports failed", error=str(e))
                if PPTX_OCR_STRICT or FORCE_OCR: raise
                return []
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY); _, bin_img = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            raw = pytesseract.image_to_string(Image.fromarray(bin_img), config="--oem 1 --psm 6")
            for l in raw.splitlines():
                if is_ocr_line_valid(l): lines.append(l.strip())
        elif OCR_BACKEND == "rapidocr":
            try:
                from rapidocr_onnxruntime import RapidOCR
            except Exception as e:
                log.warning("ocr_import_failed", "RapidOCR backend requested but import failed", error=str(e))
                if PPTX_OCR_STRICT or FORCE_OCR: raise
                return []
            ocr = RapidOCR(); res = ocr(img)
            if res and isinstance(res[0], (list, tuple)):
                for item in res[0]:
                    if len(item) >= 2:
                        text = item[1].strip()
                        if is_ocr_line_valid(text): lines.append(text)
    except Exception:
        return []
    return dedupe_lines(lines)

def import_presentation_class():
    try:
        from pptx import Presentation
        return Presentation
    except Exception:
        import importlib
        mydir = os.path.abspath(os.path.dirname(__file__))
        saved = list(sys.path)
        pruned = [p for p in saved if p and os.path.abspath(p) != mydir and os.path.abspath(p) != os.path.abspath("")]
        try:
            sys.path[:] = pruned
            mod = importlib.import_module("pptx")
            Presentation = getattr(mod, "Presentation", None)
            if Presentation is None: raise ImportError("pptx package does not expose Presentation")
            return Presentation
        finally:
            sys.path[:] = saved

def sanitize_payload_for_chunk(payload: dict) -> dict:
    try:
        range_keys = {"row_range", "slide_range", "token_range", "audio_range", "line_range"}
        for k in list(payload.keys()):
            v = payload.get(k)
            if k in range_keys and isinstance(v, (list, tuple)):
                try:
                    if len(v) == 2 and all(isinstance(x, (int, float, str)) for x in v):
                        payload[k] = [int(v[0]), int(v[1])]
                    else:
                        payload[k] = [int(x) for x in v] if all(isinstance(x, (int, float)) for x in v) else list(v)
                except Exception:
                    payload[k] = list(v)
        for list_key in ("headings", "heading_path", "tags", "layout_tags"):
            if list_key in payload:
                v = payload[list_key]
                if v is None: payload[list_key] = []
                elif isinstance(v, (list, tuple)):
                    try: payload[list_key] = [str(x) for x in v]
                    except Exception: payload[list_key] = [str(v)]
                else:
                    payload[list_key] = [str(v)]
        payload["document_id"] = str(payload.get("document_id") or "")
        payload["file_name"] = str(payload.get("file_name") or "")
        payload["chunk_id"] = str(payload.get("chunk_id") or "")
        payload["chunk_type"] = str(payload.get("chunk_type") or "")
        payload["text"] = canonicalize_text(payload.get("text") or "")
        payload["token_count"] = int(payload.get("token_count") or 0)
        payload["used_ocr"] = bool(payload.get("used_ocr", False))
        payload["layout"] = str(payload.get("layout") or "")
        if payload.get("chunk_type") == "slides" and not payload.get("slide_range"):
            payload["slide_range"] = payload.get("slide_range", [])
        return payload
    except Exception:
        safe = {
            "document_id": str(payload.get("document_id", "")),
            "file_name": str(payload.get("file_name", "")),
            "chunk_id": str(payload.get("chunk_id", "")),
            "chunk_type": str(payload.get("chunk_type", "")),
            "text": canonicalize_text(payload.get("text", "")),
            "token_count": int(payload.get("token_count", 0)),
            "tags": [],
            "headings": [],
            "heading_path": [],
            "layout_tags": [],
            "used_ocr": bool(payload.get("used_ocr", False)),
            "layout": str(payload.get("layout", "")),
        }
        return safe

def parse_file(blob_key: str, manifest: dict) -> dict:
    _init_env()
    start_all = time.perf_counter()
    client = get_storage_client_singleton()
    try:
        head_obj = client.head_object(Bucket=AZURE_CONTAINER, Key=blob_key)
    except Exception as e:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("head_failed", "Could not HEAD object %s: %s", blob_key, str(e))
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e)}
    content_len = head_obj.get("ContentLength", 0) or 0
    if content_len == 0:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.info("skip_empty", "Skipping empty object (zero bytes)", key=blob_key)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}
    try:
        raw = client.get_object(Bucket=AZURE_CONTAINER, Key=blob_key)["Body"].read()
    except Exception as e:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("read_failed", "Could not read object %s: %s", blob_key, str(e))
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e)}
    if isinstance(manifest, dict) and manifest.get("file_hash"):
        doc_id = manifest.get("file_hash")
    else:
        doc_id = sha256_hex_str(raw.decode("latin-1") if isinstance(raw, (bytes, bytearray)) else str(raw))
    out_basename = f"{doc_id}"
    raw_manifest_key = blob_key + ".manifest.json"
    try:
        if not FORCE_OVERWRITE:
            if client.exists(raw_manifest_key):
                total_ms = int((time.perf_counter() - start_all) * 1000)
                log.info("skip_manifest_exists", "raw_manifest_exists", key=raw_manifest_key)
                return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}
            if client.exists(AZURE_CHUNKED_PREFIX + out_basename + ".parquet"):
                total_ms = int((time.perf_counter() - start_all) * 1000)
                log.info("skip_parquet_exists", "parquet_exists", key=out_basename + ".parquet")
                try:
                    if not client.exists(raw_manifest_key):
                        head = client.head_object(Bucket=AZURE_CONTAINER, Key=AZURE_CHUNKED_PREFIX + out_basename + ".parquet")
                        etag = head.get("ETag", ""); etag = etag.strip('"') if isinstance(etag, str) else etag
                        size = head.get("ContentLength", 0)
                        raw_manifest = sanitize_payload_for_raw_manifest(doc_id, blob_key, AZURE_CHUNKED_PREFIX + out_basename + ".parquet", 0, etag, size)
                        client.put_object(Bucket=AZURE_CONTAINER, Key=raw_manifest_key, Body=json.dumps(raw_manifest).encode("utf-8"), ContentType="application/json")
                except Exception:
                    pass
                return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}
    except Exception:
        pass
    try:
        Presentation = import_presentation_class()
    except Exception as e:
        log.error("pptx_import_failed", "pptx import failed", error=str(e)); return {"saved_chunks": 0, "total_parse_duration_ms": 0, "skipped": True, "error": str(e)}
    try:
        prs = Presentation(BytesIO(raw))
    except Exception as e:
        log.error("pptx_open_failed", "Failed to open presentation", key=blob_key, error=str(e)); return {"saved_chunks": 0, "total_parse_duration_ms": 0, "skipped": True, "error": str(e)}
    slides_content = []
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1; t_slide_start = time.perf_counter(); text_items, table_items, img_texts = [], [], []; layout_name = ""
        try: layout_name = getattr(getattr(slide, "slide_layout", None), "name", "") or ""
        except Exception: layout_name = ""
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False):
                    txt = shape.text or ""
                    if txt.strip():
                        for ln in txt.splitlines():
                            if ln.strip(): text_items.append(ln.strip())
                if getattr(shape, "has_table", False):
                    tbl = shape.table; rows = []
                    for r in tbl.rows:
                        cols = []
                        for c in r.cells: cols.append((c.text or "").replace("\n", " ").strip())
                        rows.append(cols)
                    norm = [[cell for cell in row] for row in rows]
                    if len(norm) >= 2:
                        header = "| " + " | ".join(norm[0]) + " |"; sep = "| " + " | ".join(["---"] * len(norm[0])) + " |"
                        rows_md = ["| " + " | ".join(r) + " |" for r in norm[1:]] if len(norm) > 1 else ["\t".join(r) for r in norm]
                        md_table = "\n".join([header, sep] + rows_md) if len(norm) > 1 else "\n".join(rows_md); table_items.append(md_table)
                blob = _extract_image_blob_from_shape(shape)
                if blob and len(blob) >= MIN_IMG_BYTES:
                    try:
                        from PIL import Image
                        img = Image.open(BytesIO(blob)).convert("RGB"); arr = None
                        if _np is not None: arr = _np.array(img)[:, :, ::-1]
                        if arr is not None:
                            ocr_lines = do_ocr(arr)
                            if ocr_lines: img_texts.append("\n".join(ocr_lines))
                    except Exception:
                        pass
            except Exception:
                continue
        merged_lines = []
        if text_items: merged_lines.extend(text_items)
        if table_items: merged_lines.extend(table_items)
        if img_texts: merged_lines.extend(img_texts)
        merged_lines = [ln for ln in merged_lines if is_ocr_line_valid(ln)]; merged_lines = dedupe_lines(merged_lines)
        slide_parse_ms = (time.perf_counter() - t_slide_start) * 1000.0
        slides_content.append({"slide_number": slide_num, "raw_lines": merged_lines, "has_text": bool(text_items), "has_images_text": bool(img_texts), "tables": table_items, "parse_duration_ms": slide_parse_ms, "layout": layout_name or ""})
    saved = 0; total_slides = len(slides_content); writer = ParquetWriter(doc_id=doc_id)
    try:
        for i in range(0, total_slides, SLIDES_PER_CHUNK):
            chunk_slides = slides_content[i:i + SLIDES_PER_CHUNK]
            start = chunk_slides[0]["slide_number"]; end = chunk_slides[-1]["slide_number"]
            chunk_id = f"{doc_id}_slides_{start}_{end}"; t_chunk_start = time.perf_counter(); merged, used_ocr, slides_sum_ms, layouts = [], False, 0.0, []
            for slide in chunk_slides:
                merged.append(f"## Slide {slide['slide_number']}")
                for ln in slide["raw_lines"]: merged.append(ln)
                if slide["tables"]: merged.extend(slide["tables"])
                if slide["has_images_text"]: used_ocr = True
                if not slide["has_text"] and slide["has_images_text"]: used_ocr = True
                layouts.append(str(slide.get("layout", "") or "")); slides_sum_ms += float(slide.get("parse_duration_ms", 0.0))
            clean = [ln for ln in merged if is_ocr_line_valid(ln)]; clean = dedupe_lines(clean); final_text = "\n\n".join(clean)
            token_count = _count_tokens(final_text); merge_write_ms = (time.perf_counter() - t_chunk_start) * 1000.0; duration_ms = int(slides_sum_ms + merge_write_ms)
            dedup_layouts = []
            for l in layouts:
                if l and l not in dedup_layouts: dedup_layouts.append(l)
            layout_str = ";".join(dedup_layouts) if dedup_layouts else ""
            payload = {
                "document_id": doc_id or "",
                "file_name": os.path.basename(blob_key),
                "chunk_id": chunk_id or "",
                "chunk_type": "slides",
                "text": final_text or "",
                "token_count": int(token_count or 0),
                "figures": [],
                "embedding": None,
                "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "source_url": f"az://{AZURE_CONTAINER}/{blob_key}",
                "slide_range": [int(start), int(end)],
                "timestamp": (datetime.utcfromtimestamp(int(os.getenv('SOURCE_DATE_EPOCH'))).isoformat() + "Z") if os.getenv("SOURCE_DATE_EPOCH") else (datetime.utcnow().isoformat() + "Z"),
                "parser_version": PARSER_VERSION_PPTX,
                "tags": manifest.get("tags", []) if isinstance(manifest, dict) else [],
                "layout_tags": [],
                "layout": layout_str,
                "used_ocr": bool(used_ocr),
                "heading_path": [],
                "headings": [],
                "line_range": None
            }
            payload = sanitize_payload_for_chunk(payload)
            writer.write_payload(payload)
            log.info("buffered_chunk", "Buffered slides %d-%d (tokens=%d)", start, end, token_count)
            saved += 1
    except Exception as e:
        log.exception("buffering_failed", "Fatal error while buffering chunks", error=str(e))
        return {"saved_chunks": 0, "total_parse_duration_ms": int((time.perf_counter() - start_all) * 1000), "skipped": True, "error": str(e)}
    try:
        if saved == 0:
            total_ms = int((time.perf_counter() - start_all) * 1000); log.info("no_chunks", "No chunks produced", key=blob_key)
            return {"saved_chunks": 0, "total_parse_duration_ms": total_ms}
        count, uploaded_key, sha, size = writer.finalize_and_upload(out_basename); total_ms = int((time.perf_counter() - start_all) * 1000)
        try:
            raw_manifest = sanitize_payload_for_raw_manifest(doc_id, blob_key, uploaded_key, count, sha, size)
            client.put_object(Bucket=AZURE_CONTAINER, Key=raw_manifest_key, Body=json.dumps(raw_manifest).encode("utf-8"), ContentType="application/json")
        except Exception:
            log.warning("manifest_write_failed", "Failed to write raw manifest", key=blob_key)
        log.info("write_complete", "Wrote %d chunks", count=count, raw=blob_key, chunked=uploaded_key, duration_ms=total_ms)
        return {"saved_chunks": count, "total_parse_duration_ms": total_ms, "skipped": False}
    except Exception as e_up:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("upload_failed", "Failed to upload chunked file", key=blob_key, error=str(e_up))
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e_up)}

if __name__ == "__main__":
    _init_env()
    log.info("startup", "pptx parser start", container=AZURE_CONTAINER, env=ENV, use_managed_identity=str(USE_MANAGED_IDENTITY).lower())
    client = get_storage_client_singleton()
    paginator = client.get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=AZURE_CONTAINER, Prefix=AZURE_RAW_PREFIX):
        for obj in page.get("Contents", []):
            key = obj["Key"]
            if not key.lower().endswith(".pptx"): continue
            log.info("cli_route", "Routing parse_file", key=key)
            manifest_key = key + ".manifest.json"
            try:
                mf_obj = client.get_object(Bucket=AZURE_CONTAINER, Key=manifest_key); manifest = json.load(mf_obj["Body"])
            except Exception:
                manifest = {}
            try:
                result = parse_file(key, manifest); log.info("cli_result", "Result for key", key=key, result=result)
            except Exception:
                log.exception("cli_parse_failed", "Failed to parse", key=key)
