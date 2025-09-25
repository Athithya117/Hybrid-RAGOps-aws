#!/usr/bin/env python3
import os
import sys
import json
import logging
import boto3
import numpy as np
import time
import hashlib
import tempfile
from io import BytesIO
from datetime import datetime
from typing import Tuple, Optional

try:
    import colorama
    colorama.init()
except Exception:
    pass

RESET = "\033[0m"
COLORS = {
    logging.DEBUG: "\033[90m",
    logging.INFO: "\033[37m",
    logging.WARNING: "\033[33m",
    logging.ERROR: "\033[31m",
    logging.CRITICAL: "\033[1;41m"
}

class ColorFormatter(logging.Formatter):
    def format(self, record):
        color = COLORS.get(record.levelno, RESET)
        message = super().format(record)
        return f"{color}{message}{RESET}"

log = logging.getLogger("pptx_parser")
log.setLevel(logging.INFO)
handler = logging.StreamHandler(sys.stdout)
handler.setFormatter(ColorFormatter("%(asctime)s %(levelname)s %(message)s"))
log.handlers[:] = [handler]

REQUIRED = [
    "S3_BUCKET", "S3_RAW_PREFIX", "S3_CHUNKED_PREFIX",
    "CHUNK_FORMAT", "PPTX_SLIDES_PER_CHUNK", "PPTX_OCR_ENGINE"
]
missing = [v for v in REQUIRED if os.getenv(v) is None]
if missing:
    sys.exit(f"ERROR: Missing env vars: {', '.join(missing)}")

S3_BUCKET = os.getenv("S3_BUCKET")
S3_RAW_PREFIX = os.getenv("S3_RAW_PREFIX").rstrip("/") + "/"
S3_CHUNKED_PREFIX = os.getenv("S3_CHUNKED_PREFIX").rstrip("/") + "/"
CHUNK_FORMAT = os.getenv("CHUNK_FORMAT", "json").lower()
SLIDES_PER_CHUNK = int(os.getenv("PPTX_SLIDES_PER_CHUNK", "3"))
DISABLE_OCR = os.getenv("PPTX_DISABLE_OCR", "false").lower() == "true"
FORCE_OCR = os.getenv("PPTX_FORCE_OCR", "false").lower() == "true"
OCR_BACKEND = os.getenv("PPTX_OCR_ENGINE", "tesseract").lower()
MIN_IMG_BYTES = int(os.getenv("PPTX_MIN_IMG_SIZE_BYTES", "3072"))
PARSER_VERSION_PPTX = os.getenv("PARSER_VERSION_PPTX", "pptx-parser-v1")
TOKEN_ENCODER = os.getenv("TOKEN_ENCODER", "cl100k_base")
FORCE_OVERWRITE = os.getenv("FORCE_OVERWRITE", "false").lower() == "true"
assert CHUNK_FORMAT in ("json", "jsonl")

s3 = boto3.client("s3")


def sha256_hex(s: str) -> str:
    return hashlib.sha256((s or "").encode("utf-8")).hexdigest()


def is_valid_text(text: str) -> bool:
    t = (text or "").strip()
    return len(t) > 20 and any(c.isalpha() for c in t)


def is_ocr_line_valid(text: str, min_ratio: float = 0.6) -> bool:
    t = (text or "").strip()
    if len(t) < 5:
        return False
    alnum = sum(c.isalnum() for c in t)
    return (alnum / len(t)) >= min_ratio


def dedupe_lines(lines: list) -> list:
    seen, out = set(), []
    for l in lines:
        key = (l or "").strip().lower()
        if key and key not in seen:
            seen.add(key)
            out.append(l)
    return out


def do_ocr(img: np.ndarray) -> list:
    lines = []
    try:
        if OCR_BACKEND == "tesseract":
            import cv2
            from PIL import Image
            import pytesseract
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            _, bin_img = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            raw = pytesseract.image_to_string(Image.fromarray(bin_img), config="--oem 1 --psm 6")
            for l in raw.splitlines():
                if is_ocr_line_valid(l):
                    lines.append(l.strip())
        elif OCR_BACKEND == "rapidocr":
            from rapidocr_onnxruntime import RapidOCR
            ocr = RapidOCR()
            res = ocr(img)
            if res and isinstance(res[0], (list, tuple)):
                for item in res[0]:
                    if len(item) >= 2:
                        text = item[1].strip()
                        if is_ocr_line_valid(text):
                            lines.append(text)
    except Exception:
        return []
    return dedupe_lines(lines)


def is_valid_table(table: list) -> bool:
    if len(table) < 2 or len(table[0]) < 2:
        return False
    total_cells = sum(len(r) for r in table)
    alpha_cells = sum(1 for row in table for cell in row if any(c.isalpha() for c in (cell or "")))
    return total_cells > 0 and (alpha_cells / total_cells) >= 0.5


def _extract_image_blob_from_shape(shape):
    try:
        img = getattr(shape, "image", None)
        if img and getattr(img, "blob", None):
            return img.blob
    except Exception:
        pass
    try:
        fill = getattr(shape, "fill", None)
        if fill is not None and getattr(fill, "type", None) is not None:
            pic = getattr(fill, "picture", None)
            if pic and getattr(pic, "image", None) and getattr(pic.image, "blob", None):
                return pic.image.blob
    except Exception:
        pass
    return None


def _count_tokens(text: str) -> int:
    if not text:
        return 0
    try:
        import tiktoken
        enc = tiktoken.get_encoding(TOKEN_ENCODER)
        return len(enc.encode(text))
    except Exception:
        return len(text.split())


def s3_object_exists(key: str) -> bool:
    try:
        s3.head_object(Bucket=S3_BUCKET, Key=key)
        return True
    except Exception:
        return False


class S3DocWriter:
    """
    Aggregate chunks for a single source into one temp file in /tmp and upload once.
    Ensures proper formatting:
      - jsonl: one JSON object per line
      - json: pretty array where each object is indented and separated with comma+newline
    """
    def __init__(self, doc_id: str, s3_path: str, ext: str, content_type: str = "application/json"):
        self.doc_id = doc_id
        self.s3_path = s3_path or ""
        self.ext = ext
        self.content_type = content_type
        # force temp files into /tmp to avoid filling other dirs; delete=False so we can upload then remove
        self.temp = tempfile.NamedTemporaryFile(mode="wb", delete=False, suffix=f".{ext}", dir="/tmp")
        self.count = 0
        self._first = True
        if self.ext == "json":
            # start pretty array with newline to allow multi-line objects per element
            self.temp.write(b"[\n")
            self.temp.flush()

    def write_payload(self, payload: dict) -> int:
        self.count += 1
        if self.ext == "jsonl":
            line = (json.dumps(payload, ensure_ascii=False) + "\n").encode("utf-8")
            self.temp.write(line)
        else:
            pretty = json.dumps(payload, ensure_ascii=False, indent=2)
            # indent each line by two spaces for readability
            indented = ("\n".join("  " + ln for ln in pretty.splitlines()) + "\n").encode("utf-8")
            if not self._first:
                self.temp.write(b",\n")
            self.temp.write(indented)
            self._first = False
        self.temp.flush()
        return 1

    def finalize_and_upload(self, out_key: str) -> Tuple[int, str]:
        if self.ext == "json":
            self.temp.write(b"]\n")
        self.temp.flush()
        self.temp.close()
        try:
            s3.upload_file(self.temp.name, S3_BUCKET, out_key, ExtraArgs={"ContentType": self.content_type})
            try:
                os.unlink(self.temp.name)
            except Exception:
                pass
            return self.count, out_key
        except Exception:
            try:
                os.unlink(self.temp.name)
            except Exception:
                pass
            raise


def _derive_doc_id_from_head(s3_key: str, head_obj: dict, manifest: dict) -> str:
    """
    Derive a stable doc_id without downloading content.
    Priority:
      1. manifest['file_hash'] (if present)
      2. ETag (from HEAD) - stripped of quotes
      3. LastModified (string)
      4. fallback to basename(s3_key)
      5. fallback to sha256(s3_key)
    """
    if isinstance(manifest, dict) and manifest.get("file_hash"):
        return manifest.get("file_hash")
    etag = head_obj.get("ETag", "")
    if isinstance(etag, str):
        etag = etag.strip('"')
    if etag:
        return sha256_hex(s3_key + str(etag))
    lm = head_obj.get("LastModified", "")
    if lm:
        return sha256_hex(s3_key + str(lm))
    base = os.path.basename(s3_key)
    if base:
        return base
    return sha256_hex(s3_key)


def parse_file(s3_key: str, manifest: dict) -> dict:
    start_all = time.perf_counter()

    # 1) HEAD for quick metadata and doc_id derivation
    try:
        head_obj = s3.head_object(Bucket=S3_BUCKET, Key=s3_key)
    except Exception as e:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("Could not HEAD S3 object %s: %s", s3_key, e)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e)}

    content_len = head_obj.get("ContentLength", 0) or 0

    # derive doc_id and out_key without downloading
    doc_id = _derive_doc_id_from_head(s3_key, head_obj, manifest or {})
    ext = "jsonl" if CHUNK_FORMAT == "jsonl" else "json"
    out_key = f"{S3_CHUNKED_PREFIX}{doc_id}.{ext}"

    # fast-skip if aggregated output exists
    if not FORCE_OVERWRITE and s3_object_exists(out_key):
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.info("Skipping entire file because chunked file exists: %s", out_key)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}

    # skip empty files
    if content_len == 0:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.info("Skipping empty object %s (zero bytes).", s3_key)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}

    # now GET full object (we need it)
    try:
        raw = s3.get_object(Bucket=S3_BUCKET, Key=s3_key)["Body"].read()
    except Exception as e:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.error("Could not read S3 object %s: %s", s3_key, e)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e)}

    # if manifest.file_hash exists prefer that (keeps stable previous behavior)
    if isinstance(manifest, dict) and manifest.get("file_hash"):
        doc_id = manifest.get("file_hash")
        out_key = f"{S3_CHUNKED_PREFIX}{doc_id}.{ext}"

    # race-check after download (someone may have written the aggregated file between HEAD and GET)
    if not FORCE_OVERWRITE and s3_object_exists(out_key):
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.info("Skipping entire file because chunked file exists (post-download): %s", out_key)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}

    # ensure pptx library available
    try:
        from pptx import Presentation
    except Exception as e:
        log.error(f"pptx import failed: {e}")
        return {"saved_chunks": 0, "total_parse_duration_ms": 0, "skipped": True, "error": str(e)}

    # parse presentation
    try:
        prs = Presentation(BytesIO(raw))
    except Exception as e:
        log.error("Failed to open presentation %s: %s", s3_key, e)
        return {"saved_chunks": 0, "total_parse_duration_ms": 0, "skipped": True, "error": str(e)}

    slides_content = []
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        t_slide_start = time.perf_counter()
        text_items = []
        table_items = []
        img_texts = []
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False):
                    txt = shape.text or ""
                    if txt.strip():
                        for ln in txt.splitlines():
                            if ln.strip():
                                text_items.append(ln.strip())
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    rows = []
                    for r in tbl.rows:
                        cols = []
                        for c in r.cells:
                            cols.append((c.text or "").replace("\n", " ").strip())
                        rows.append(cols)
                    norm = [[cell for cell in row] for row in rows]
                    if is_valid_table(norm):
                        header = "| " + " | ".join(norm[0]) + " |"
                        sep = "| " + " | ".join(["---"] * len(norm[0])) + " |"
                        rows_md = ["| " + " | ".join(r) + " |" for r in norm[1:]] if len(norm) > 1 else ["\t".join(r) for r in norm]
                        md_table = "\n".join([header, sep] + rows_md) if len(norm) > 1 else "\n".join(rows_md)
                        table_items.append(md_table)
                blob = _extract_image_blob_from_shape(shape)
                if blob and len(blob) >= MIN_IMG_BYTES:
                    from PIL import Image
                    img = Image.open(BytesIO(blob)).convert("RGB")
                    arr = np.array(img)[:, :, ::-1]
                    ocr_lines = do_ocr(arr)
                    if ocr_lines:
                        img_texts.append("\n".join(ocr_lines))
            except Exception:
                # ignore shape-level failures and continue
                continue
        merged_lines = []
        if text_items:
            merged_lines.extend(text_items)
        if table_items:
            merged_lines.extend(table_items)
        if img_texts:
            merged_lines.extend(img_texts)
        merged_lines = [ln for ln in merged_lines if is_ocr_line_valid(ln)]
        merged_lines = dedupe_lines(merged_lines)
        slide_parse_ms = (time.perf_counter() - t_slide_start) * 1000.0
        slides_content.append({
            "slide_number": slide_num,
            "raw_lines": merged_lines,
            "has_text": bool(text_items),
            "has_images_text": bool(img_texts),
            "tables": table_items,
            "parse_duration_ms": slide_parse_ms
        })

    saved = 0
    total_slides = len(slides_content)
    ext = "jsonl" if CHUNK_FORMAT == "jsonl" else "json"
    out_key = f"{S3_CHUNKED_PREFIX}{doc_id}.{ext}"

    # final skip check (defensive)
    if not FORCE_OVERWRITE and s3_object_exists(out_key):
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.info("Skipping entire file because chunked file exists: %s", out_key)
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True}

    writer = S3DocWriter(doc_id=doc_id, s3_path=s3_key, ext=ext)
    try:

        def _sanitize_payload_for_weaviate(payload: dict) -> dict:
            """
            Weaviate often defines range-like fields as TEXT. If a parser emits lists like [start,end]
            we should convert them to a string to avoid insertion errors.
            This helper:
              - converts known range fields (row_range, slide_range, token_range, audio_range, line_range)
                into a short "start-end" string when appropriate, otherwise JSON stringifies them.
              - ensures headings, heading_path and tags are stringified elements (keeps them as arrays).
            """
            range_keys = {"row_range", "slide_range", "token_range", "audio_range", "line_range"}
            for k in list(payload.keys()):
                v = payload.get(k)
                if k in range_keys and isinstance(v, (list, tuple)):
                    # prefer "start-end" for simple pair ranges, otherwise json-dump
                    try:
                        if len(v) == 2 and all(isinstance(x, (int, str)) for x in v):
                            payload[k] = f"{v[0]}-{v[1]}"
                        else:
                            payload[k] = json.dumps(v)
                    except Exception:
                        payload[k] = json.dumps(v)
            # ensure arrays of headings/tags are arrays of strings (Weaviate TEXT_ARRAY expects strings)
            if "headings" in payload and isinstance(payload["headings"], (list, tuple)):
                payload["headings"] = [str(x) for x in payload["headings"]]
            if "heading_path" in payload and isinstance(payload["heading_path"], (list, tuple)):
                payload["heading_path"] = [str(x) for x in payload["heading_path"]]
            if "tags" in payload and isinstance(payload["tags"], (list, tuple)):
                payload["tags"] = [str(x) for x in payload["tags"]]
            return payload

        for i in range(0, total_slides, SLIDES_PER_CHUNK):
            chunk_slides = slides_content[i:i + SLIDES_PER_CHUNK]
            start = chunk_slides[0]["slide_number"]
            end = chunk_slides[-1]["slide_number"]
            chunk_id = f"{doc_id}_slides_{start}_{end}"
            t_chunk_start = time.perf_counter()
            merged = []
            used_ocr = False
            slides_sum_ms = 0.0
            for slide in chunk_slides:
                merged.append(f"## Slide {slide['slide_number']}")
                for ln in slide["raw_lines"]:
                    merged.append(ln)
                if slide["tables"]:
                    merged.extend(slide["tables"])
                if slide["has_images_text"]:
                    used_ocr = True
                if not slide["has_text"] and slide["has_images_text"]:
                    used_ocr = True
                slides_sum_ms += float(slide.get("parse_duration_ms", 0.0))
            clean = [ln for ln in merged if is_ocr_line_valid(ln)]
            clean = dedupe_lines(clean)
            final_text = "\n\n".join(clean)
            token_count = _count_tokens(final_text)
            merge_write_ms = (time.perf_counter() - t_chunk_start) * 1000.0
            duration_ms = int(slides_sum_ms + merge_write_ms)
            payload = {
                "document_id": doc_id or "",
                "file_name": os.path.basename(s3_key),
                "chunk_id": chunk_id or "",
                "chunk_type": "slides",
                "text": final_text or "",
                "token_count": int(token_count or 0),
                "embedding": None,
                "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                "source_url": f"s3://{S3_BUCKET}/{s3_key}",
                "page_number": None,
                "slide_range": [int(start), int(end)],
                "row_range": None,
                "token_range": None,
                "audio_range": None,
                "timestamp": datetime.utcnow().isoformat() + "Z",
                "parser_version": PARSER_VERSION_PPTX,
                "tags": manifest.get("tags", []) if isinstance(manifest, dict) else [],
                "layout_tags": ["slide"],
                "used_ocr": bool(used_ocr),
                "heading_path": [],
                "headings": [],
                "line_range": None
            }
            # sanitize problematic non-string fields (ranges) so Weaviate doesn't reject them
            payload = _sanitize_payload_for_weaviate(payload)
            writer.write_payload(payload)
            log.info("Buffered slides %d-%d (tokens=%d)", start, end, token_count)
            saved += 1
    except Exception as e:
        try:
            if writer and writer.temp:
                try:
                    os.unlink(writer.temp.name)
                except Exception:
                    pass
        except Exception:
            pass
        log.exception("Fatal error while buffering chunks for %s: %s", s3_key, str(e))
        return {"saved_chunks": 0, "total_parse_duration_ms": int((time.perf_counter() - start_all) * 1000), "skipped": True, "error": str(e)}

    try:
        if saved == 0:
            try:
                if writer and writer.temp:
                    os.unlink(writer.temp.name)
            except Exception:
                pass
            total_ms = int((time.perf_counter() - start_all) * 1000)
            log.info("No chunks produced for %s", s3_key)
            return {"saved_chunks": 0, "total_parse_duration_ms": total_ms}
        count, uploaded_key = writer.finalize_and_upload(out_key)
        total_ms = int((time.perf_counter() - start_all) * 1000)
        log.info("Wrote %d chunks for %s â†’ %s (%d ms total)", count, s3_key, uploaded_key, total_ms)
        return {"saved_chunks": count, "total_parse_duration_ms": total_ms}
    except Exception as e_up:
        total_ms = int((time.perf_counter() - start_all) * 1000)
        try:
            if writer and writer.temp:
                try:
                    os.unlink(writer.temp.name)
                except Exception:
                    pass
        except Exception:
            pass
        log.error("Failed to upload chunked file for %s error=%s", s3_key, str(e_up))
        return {"saved_chunks": 0, "total_parse_duration_ms": total_ms, "skipped": True, "error": str(e_up)}


if __name__ == "__main__":
    paginator = s3.get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=S3_RAW_PREFIX):
        for obj in page.get("Contents", []):
            key = obj["Key"]
            if not key.lower().endswith(".pptx"):
                continue
            log.info("Routing parse_file for s3://%s/%s", S3_BUCKET, key)
            manifest_key = key + ".manifest.json"
            try:
                mf = s3.get_object(Bucket=S3_BUCKET, Key=manifest_key)
                manifest = json.load(mf["Body"])
            except Exception:
                manifest = {}
            parse_file(key, manifest)
