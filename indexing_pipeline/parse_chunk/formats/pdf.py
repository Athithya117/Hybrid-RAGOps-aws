#!/usr/bin/env python3
import os
import sys
import json
import logging
import boto3
import numpy as np
import time
import hashlib
from io import BytesIO
from datetime import datetime
from botocore.exceptions import ClientError
from tempfile import NamedTemporaryFile
from typing import Optional

try:
    import colorama
    colorama.init()
except Exception:
    pass

RESET = "\033[0m"
COLORS = {
    logging.DEBUG: "\033[90m",
    logging.INFO: "\033[37m",
    logging.WARNING: "\033[33m",
    logging.ERROR: "\033[31m",
    logging.CRITICAL: "\033[1;41m",
}


class ColorFormatter(logging.Formatter):
    def format(self, record):
        color = COLORS.get(record.levelno, RESET)
        message = super().format(record)
        return f"{color}{message}{RESET}"


log = logging.getLogger("pptx_parser")
log.setLevel(logging.INFO)
handler = logging.StreamHandler(sys.stdout)
handler.setFormatter(ColorFormatter("%(asctime)s %(levelname)s %(message)s"))
log.handlers[:] = [handler]

REQUIRED = [
    "S3_BUCKET",
    "S3_RAW_PREFIX",
    "S3_CHUNKED_PREFIX",
    "CHUNK_FORMAT",
    "PPTX_SLIDES_PER_CHUNK",
    "PPTX_OCR_ENGINE",
]
missing = [v for v in REQUIRED if os.getenv(v) is None]
if missing:
    sys.exit(f"ERROR: Missing env vars: {', '.join(missing)}")

S3_BUCKET = os.getenv("S3_BUCKET")
S3_RAW_PREFIX = os.getenv("S3_RAW_PREFIX").rstrip("/") + "/"
S3_CHUNKED_PREFIX = os.getenv("S3_CHUNKED_PREFIX").rstrip("/") + "/"
CHUNK_FORMAT = os.getenv("CHUNK_FORMAT", "json").lower()
SLIDES_PER_CHUNK = int(os.getenv("PPTX_SLIDES_PER_CHUNK", "3"))
DISABLE_OCR = os.getenv("PPTX_DISABLE_OCR", "false").lower() == "true"
FORCE_OCR = os.getenv("PPTX_FORCE_OCR", "false").lower() == "true"
OCR_BACKEND = os.getenv("PPTX_OCR_ENGINE", "tesseract").lower()
MIN_IMG_BYTES = int(os.getenv("PPTX_MIN_IMG_SIZE_BYTES", "3072"))
PARSER_VERSION_PPTX = os.getenv("PARSER_VERSION_PPTX", "py-pptx-v1")
TOKEN_ENCODER = os.getenv("TOKEN_ENCODER", "cl100k_base")
STORE_ONE_FILE_PER_CHUNK = os.getenv("STORE_ONE_FILE_PER_CHUNK", "false").lower() == "true"
assert CHUNK_FORMAT in ("json", "jsonl")

s3 = boto3.client("s3")


def sha256_hex_str_of_bytes(b: bytes) -> str:
    h = hashlib.sha256()
    h.update(b)
    return h.hexdigest()


def sha256_hex(s: str) -> str:
    return hashlib.sha256((s or "").encode("utf-8")).hexdigest()


def is_valid_text(text: str) -> bool:
    t = (text or "").strip()
    return len(t) > 20 and any(c.isalpha() for c in t)


def is_ocr_line_valid(text: str, min_ratio: float = 0.6) -> bool:
    t = (text or "").strip()
    if len(t) < 5:
        return False
    alnum = sum(c.isalnum() for c in t)
    return (alnum / len(t)) >= min_ratio


def dedupe_lines(lines: list[str]) -> list[str]:
    seen, out = set(), []
    for l in lines:
        key = (l or "").strip().lower()
        if key and key not in seen:
            seen.add(key)
            out.append(l)
    return out


def do_ocr(img: np.ndarray) -> list[str]:
    lines = []
    try:
        if OCR_BACKEND == "tesseract":
            import cv2
            from PIL import Image
            import pytesseract
            gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
            _, bin_img = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
            raw = pytesseract.image_to_string(Image.fromarray(bin_img), config="--oem 1 --psm 6")
            for l in raw.splitlines():
                if is_ocr_line_valid(l):
                    lines.append(l.strip())
        elif OCR_BACKEND == "rapidocr":
            from rapidocr_onnxruntime import RapidOCR
            ocr = RapidOCR()
            res = ocr(img)
            if res and isinstance(res[0], (list, tuple)):
                for item in res[0]:
                    if len(item) >= 2:
                        text = item[1].strip()
                        if is_ocr_line_valid(text):
                            lines.append(text)
    except Exception:
        return []
    return dedupe_lines(lines)


def is_valid_table(table: list[list[str]]) -> bool:
    if len(table) < 2 or len(table[0]) < 2:
        return False
    total_cells = sum(len(r) for r in table)
    alpha_cells = sum(1 for row in table for cell in row if any(c.isalpha() for c in (cell or "")))
    return total_cells > 0 and (alpha_cells / total_cells) >= 0.5


def _extract_image_blob_from_shape(shape):
    try:
        img = getattr(shape, "image", None)
        if img and getattr(img, "blob", None):
            return img.blob
    except Exception:
        pass
    try:
        fill = getattr(shape, "fill", None)
        if fill is not None and getattr(fill, "type", None) is not None:
            pic = getattr(fill, "picture", None)
            if pic and getattr(pic, "image", None) and getattr(pic.image, "blob", None):
                return pic.image.blob
    except Exception:
        pass
    return None


def _count_tokens(text: str) -> int:
    if not text:
        return 0
    try:
        import tiktoken
        enc = tiktoken.get_encoding(TOKEN_ENCODER)
        return len(enc.encode(text))
    except Exception:
        return len(text.split())


def _s3_object_metadata_sha256(bucket: str, key: str) -> Optional[str]:
    try:
        ho = s3.head_object(Bucket=bucket, Key=key)
        meta = ho.get("Metadata", {}) or {}
        return meta.get("sha256")
    except ClientError:
        return None


class LocalChunkAppender:
    def __init__(self, chunk_format: str, doc_id: str):
        self.chunk_format = chunk_format
        self.doc_id = doc_id
        suffix = f".{chunk_format}"
        self.temp = NamedTemporaryFile(delete=False, mode="w", encoding="utf-8", suffix=suffix)
        self.path = self.temp.name
        self.count = 0
    def append(self, payload: dict):
        if self.chunk_format == "jsonl":
            line = json.dumps(payload, ensure_ascii=False, sort_keys=True)
            self.temp.write(line + "\n")
        else:
            pretty = json.dumps(payload, indent=2, ensure_ascii=False, sort_keys=True)
            self.temp.write(pretty + "\n")
        self.count += 1
        self.temp.flush()
    def finalize_and_upload(self, s3_bucket: str, s3_key: str):
        self.temp.close()
        with open(self.path, "rb") as f:
            body_bytes = f.read()
        body_sha = sha256_hex_str_of_bytes(body_bytes)
        existing_sha = _s3_object_metadata_sha256(s3_bucket, s3_key)
        if existing_sha and existing_sha == body_sha:
            try:
                os.remove(self.path)
            except Exception:
                pass
            log.info(f"Skipping upload for {self.doc_id} → s3://{s3_bucket}/{s3_key} (identical)")
            return
        extra = {"ContentType": "application/json", "Metadata": {"sha256": body_sha}}
        try:
            s3.upload_file(self.path, s3_bucket, s3_key, ExtraArgs=extra)
            log.info(f"Uploaded combined chunks for {self.doc_id} → s3://{s3_bucket}/{s3_key} ({self.count} chunks)")
        finally:
            try:
                os.remove(self.path)
            except Exception:
                pass


def parse_file(s3_key: str, manifest: dict) -> dict:
    start_all = time.perf_counter()
    try:
        raw = s3.get_object(Bucket=S3_BUCKET, Key=s3_key)["Body"].read()
    except Exception as e:
        log.error(f"Failed to read {s3_key} from S3: {e}")
        return {"saved_chunks": 0, "total_parse_duration_ms": 0}
    doc_id = manifest.get("file_hash") if isinstance(manifest, dict) else None
    if not doc_id:
        doc_id = os.path.basename(s3_key) or sha256_hex(s3_key)
    source = f"s3://{S3_BUCKET}/{s3_key}"
    try:
        from pptx import Presentation
    except Exception as e:
        log.error(f"pptx import failed: {e}")
        return {"saved_chunks": 0, "total_parse_duration_ms": 0}
    try:
        prs = Presentation(BytesIO(raw))
    except Exception as e:
        log.error(f"Failed to parse PPTX {s3_key}: {e}")
        return {"saved_chunks": 0, "total_parse_duration_ms": 0}
    slides_content = []
    for idx, slide in enumerate(prs.slides):
        slide_num = idx + 1
        t_slide_start = time.perf_counter()
        text_items = []
        table_items = []
        img_texts = []
        for shape in slide.shapes:
            try:
                if getattr(shape, "has_text_frame", False):
                    txt = shape.text or ""
                    if txt.strip():
                        for ln in txt.splitlines():
                            if ln.strip():
                                text_items.append(ln.strip())
                if getattr(shape, "has_table", False):
                    tbl = shape.table
                    rows = []
                    for r in tbl.rows:
                        cols = []
                        for c in r.cells:
                            cols.append((c.text or "").replace("\n", " ").strip())
                        rows.append(cols)
                    norm = [[cell for cell in row] for row in rows]
                    if is_valid_table(norm):
                        header = "| " + " | ".join(norm[0]) + " |"
                        sep = "| " + " | ".join(["---"] * len(norm[0])) + " |"
                        rows_md = ["| " + " | ".join(r) + " |" for r in norm[1:]] if len(norm) > 1 else ["\t".join(r) for r in norm]
                        md_table = "\n".join([header, sep] + rows_md) if len(norm) > 1 else "\n".join(rows_md)
                        table_items.append(md_table)
                blob = _extract_image_blob_from_shape(shape)
                if blob and len(blob) >= MIN_IMG_BYTES:
                    from PIL import Image
                    img = Image.open(BytesIO(blob)).convert("RGB")
                    arr = np.array(img)[:, :, ::-1]
                    ocr_lines = do_ocr(arr)
                    if ocr_lines:
                        img_texts.append("\n".join(ocr_lines))
            except Exception:
                continue
        merged_lines = []
        if text_items:
            merged_lines.extend(text_items)
        if table_items:
            merged_lines.extend(table_items)
        if img_texts:
            merged_lines.extend(img_texts)
        merged_lines = [ln for ln in merged_lines if is_ocr_line_valid(ln)]
        merged_lines = dedupe_lines(merged_lines)
        slide_parse_ms = (time.perf_counter() - t_slide_start) * 1000.0
        slides_content.append({
            "slide_number": slide_num,
            "raw_lines": merged_lines,
            "has_text": bool(text_items),
            "has_images_text": bool(img_texts),
            "tables": table_items,
            "parse_duration_ms": slide_parse_ms,
        })
    saved = 0
    total_slides = len(slides_content)
    ext = "jsonl" if CHUNK_FORMAT == "jsonl" else "json"
    combined_appender = None
    combined_key = f"{S3_CHUNKED_PREFIX}{doc_id}.{ext}"
    if not STORE_ONE_FILE_PER_CHUNK:
        combined_appender = LocalChunkAppender(ext, doc_id)
        log.info(f"Using combined chunk file mode for {doc_id} → s3://{S3_BUCKET}/{combined_key}")
    for i in range(0, total_slides, SLIDES_PER_CHUNK):
        chunk_slides = slides_content[i:i + SLIDES_PER_CHUNK]
        start = chunk_slides[0]["slide_number"]
        end = chunk_slides[-1]["slide_number"]
        chunk_id = f"{doc_id}_slides_{start}_{end}"
        t_chunk_start = time.perf_counter()
        merged = []
        used_ocr = False
        slides_sum_ms = 0.0
        for slide in chunk_slides:
            merged.append(f"## Slide {slide['slide_number']}")
            for ln in slide["raw_lines"]:
                merged.append(ln)
            if slide["tables"]:
                merged.extend(slide["tables"])
            if slide["has_images_text"]:
                used_ocr = True
            if not slide["has_text"] and slide["has_images_text"]:
                used_ocr = True
            slides_sum_ms += float(slide.get("parse_duration_ms", 0.0))
        clean = [ln for ln in merged if is_ocr_line_valid(ln)]
        clean = dedupe_lines(clean)
        final_text = "\n\n".join(clean)
        token_count = _count_tokens(final_text)
        merge_write_ms = (time.perf_counter() - t_chunk_start) * 1000.0
        duration_ms = int(slides_sum_ms + merge_write_ms)
        payload = {
            "document_id": doc_id or "",
            "chunk_id": chunk_id or "",
            "chunk_type": "slides",
            "text": final_text or "",
            "token_count": int(token_count or 0),
            "embedding": None,
            "file_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "source_url": source,
            "page_number": None,
            "slide_range": [int(start), int(end)],
            "row_range": None,
            "token_range": None,
            "audio_range": None,
            "timestamp": datetime.utcnow().isoformat() + "Z",
            "parser_version": PARSER_VERSION_PPTX,
            "tags": manifest.get("tags", []) if isinstance(manifest, dict) else [],
            "layout_tags": ["slide"],
            "used_ocr": bool(used_ocr),
            "parse_chunk_duration_ms": int(duration_ms),
            "heading_path": [],
            "headings": [],
            "line_range": None,
            "chunk_duration_ms": None,
        }
        if STORE_ONE_FILE_PER_CHUNK:
            out_key = f"{S3_CHUNKED_PREFIX}{chunk_id}.{ext}"
            if ext == "jsonl":
                body_bytes = (json.dumps(payload, ensure_ascii=False, sort_keys=True) + "\n").encode()
            else:
                body_bytes = json.dumps(payload, indent=2, ensure_ascii=False, sort_keys=True).encode()
            body_sha = sha256_hex_str_of_bytes(body_bytes)
            existing_sha = _s3_object_metadata_sha256(S3_BUCKET, out_key)
            if existing_sha and existing_sha == body_sha:
                log.info(f"Skipping upload for {out_key} (identical)")
                saved += 1
            else:
                try:
                    s3.put_object(Bucket=S3_BUCKET, Key=out_key, Body=body_bytes, ContentType="application/json", Metadata={"sha256": body_sha})
                    log.info(f"Parsed slides {start}-{end} in {duration_ms} ms (tokens={token_count}) → {out_key}")
                    saved += 1
                except ClientError as e:
                    log.error(f"Failed to write {out_key}: {e}")
        else:
            try:
                combined_appender.append(payload)
                log.info(f"Appended slides {start}-{end} in {duration_ms} ms (tokens={token_count}) → {combined_key}")
                saved += 1
            except Exception as e:
                log.error(f"Failed to append payload for {chunk_id}: {e}")
    if not STORE_ONE_FILE_PER_CHUNK and combined_appender is not None:
        try:
            combined_appender.finalize_and_upload(S3_BUCKET, combined_key)
        except Exception as e:
            log.error(f"Failed uploading combined file for {doc_id}: {e}")
    total_ms = int((time.perf_counter() - start_all) * 1000)
    log.info(f"Completed parsing {saved} chunks ({total_slides} slides) in {total_ms} ms total")
    return {"saved_chunks": saved, "total_parse_duration_ms": total_ms}


if __name__ == "__main__":
    paginator = s3.get_paginator("list_objects_v2")
    for page in paginator.paginate(Bucket=S3_BUCKET, Prefix=S3_RAW_PREFIX):
        for obj in page.get("Contents", []):
            key = obj["Key"]
            if not key.lower().endswith(".pptx"):
                continue
            log.info(f"Routing parse_file for s3://{S3_BUCKET}/{key}")
            manifest_key = key + ".manifest.json"
            try:
                mf = s3.get_object(Bucket=S3_BUCKET, Key=manifest_key)
                manifest = json.loads(mf["Body"].read().decode("utf-8"))
            except Exception:
                manifest = {}
            parse_file(key, manifest)
